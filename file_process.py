"""
file_process.py — Tải & xử lý file người dùng gửi lên (INBOUND).

Khác với file_cache.py (chỉ lo file bot GỬI RA, cache trong RAM), module này
lo file NGƯỜI DÙNG GỬI LÊN cho bot: tải về thẳng /tmp (không có thư mục
con riêng — file được đặt tên với prefix "tgbot_incoming_" để phân biệt
với file của process khác trong /tmp) bằng file_id thật của Telegram rồi
cố gắng đưa NỘI DUNG cho Gemini, thay vì chỉ ghi lại tên/kích thước như
trước.

Luồng xử lý 1 file:
  1. bot.get_file(file_id) + download_to_drive() → tải về /tmp (file thật,
     không giữ trong RAM cache).
  2. Nếu mime được Gemini đọc trực tiếp (ảnh / audio / video / PDF) → base64
     hoá thành inlineData part, đính kèm thẳng vào request — Gemini tự
     xem ảnh / nghe audio / đọc PDF (kể cả OCR), bot không cần tự parse.
  3. Nếu là .docx → trích xuất text bằng python-docx.
  4. Nếu là file text thuần (txt/md/csv/json/code...) → đọc trực tiếp.
  5. Trường hợp còn lại (không hỗ trợ, hoặc quá lớn) → fallback: chỉ giữ
     dòng metadata cũ (tên/size/file_id), không có nội dung.

File /tmp KHÔNG bị xoá ngay sau khi xử lý xong nữa — được giữ lại và dọn
bởi một vòng lặp nền (`start_cleanup_loop`, chạy trong main.py) xoá các
file cũ hơn FILE_RETENTION_SECONDS (mặc định 24h). Vì lưu thẳng trong /tmp
dùng chung cho cả container, vòng dọn CHỈ xoá file có prefix
FILENAME_PREFIX — không đụng tới file của process/thư viện khác. Lưu ý:
Render có thể restart bất kỳ lúc nào nên phần file chưa kịp dọn sẽ mất theo
container — không sao vì không có phần nào khác trong code đọc lại file
theo path này, đây chỉ là giữ lại để debug/audit trong vòng 24h.
"""
from __future__ import annotations

import asyncio
import base64
import logging
import mimetypes
import os
import time
import uuid
from dataclasses import dataclass
from typing import Optional

from telegram import Bot

from config import MAX_MEDIA_MB
from i18n import DEFAULT_LANG, t

logger = logging.getLogger(__name__)

TMP_DIR = "/tmp"

# Mọi file bot tải về đều được đặt tên với prefix này. /tmp là thư mục dùng
# chung cho CẢ container (pip, os.tempfile, các process khác...), nên vòng
# lặp dọn dẹp (_sweep_expired_once) CHỈ được phép động vào file có prefix
# này — tuyệt đối không quét/xoá toàn bộ /tmp, tránh xoá nhầm file của
# process/thư viện khác.
FILENAME_PREFIX = "tgbot_incoming_"

MAX_MEDIA_BYTES = MAX_MEDIA_MB * 1024 * 1024

# File trong TMP_DIR được giữ lại tối đa từng này giây trước khi bị dọn.
FILE_RETENTION_SECONDS = 24 * 60 * 60  # 24h

# Tần suất chạy vòng quét dọn file hết hạn.
CLEANUP_SWEEP_INTERVAL_SECONDS = 60 * 60  # 1h/lần là đủ, không cần quét dày

# MIME mà Gemini nhận trực tiếp qua inlineData — Gemini tự đọc nội dung
# (ảnh/audio/video/PDF), bot không cần tự trích xuất text.
GEMINI_INLINE_MIMES = {
    "image/jpeg", "image/png", "image/webp", "image/heic", "image/heif",
    "audio/mpeg", "audio/mp3", "audio/wav", "audio/x-wav", "audio/aac",
    "audio/ogg", "audio/flac", "audio/mp4", "audio/m4a",
    "video/mp4", "video/mpeg", "video/mov", "video/quicktime", "video/avi",
    "video/x-flv", "video/mpg", "video/webm", "video/wmv", "video/3gpp",
    "application/pdf",
}

# Đuôi file đọc thẳng như text, không cần lib ngoài.
TEXT_EXTS = {
    ".txt", ".md", ".csv", ".json", ".log", ".py", ".js", ".ts", ".tsx",
    ".jsx", ".html", ".htm", ".css", ".xml", ".yaml", ".yml", ".sql", ".sh",
    ".ini", ".cfg", ".toml", ".java", ".c", ".cpp", ".h", ".go", ".rs", ".php",
}

MAX_TEXT_CHARS = 12000  # tránh 1 file quá dài nuốt hết context budget


@dataclass
class FileResult:
    media_part: Optional[dict] = None  # part inlineData gửi thẳng cho Gemini
    extra_text: Optional[str] = None   # nội dung trích xuất (docx/text)


def _guess_mime(filename: str, tg_mime: Optional[str]) -> str:
    if tg_mime:
        return tg_mime
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or "application/octet-stream"


async def _download(bot: Bot, file_id: str, filename: str, size: Optional[int]) -> Optional[str]:
    try:
        tg_file = await bot.get_file(file_id)
    except Exception as e:
        logger.warning("[file_process] get_file thất bại (%s): %s", filename, e)
        return None

    real_size = tg_file.file_size or size or 0
    if real_size > MAX_MEDIA_BYTES:
        logger.info(
            "[file_process] '%s' quá lớn (%.1f MB > %d MB) — bỏ qua tải nội dung.",
            filename, real_size / 1024 / 1024, MAX_MEDIA_MB,
        )
        return None

    ext  = os.path.splitext(filename)[1] or ""
    path = os.path.join(TMP_DIR, f"{FILENAME_PREFIX}{uuid.uuid4().hex}{ext}")
    try:
        await tg_file.download_to_drive(path)
    except Exception as e:
        logger.warning("[file_process] download_to_drive lỗi (%s): %s", filename, e)
        return None
    return path


def _sweep_expired_once(now: Optional[float] = None) -> int:
    """Quét TMP_DIR 1 lần, xoá file CỦA BOT (prefix FILENAME_PREFIX) có
    mtime cũ hơn FILE_RETENTION_SECONDS.

    TMP_DIR = "/tmp" dùng chung cho cả container nên hàm này CHỈ được xét
    entry.name.startswith(FILENAME_PREFIX) — không bao giờ xoá file không
    phải do file_process.py tạo ra.

    Trả về số file đã xoá. Chạy đồng bộ (os.scandir/os.remove nhanh, không
    cần async) — được gọi định kỳ bởi start_cleanup_loop().
    """
    now = now if now is not None else time.time()
    removed = 0
    try:
        entries = os.scandir(TMP_DIR)
    except OSError as e:
        logger.warning("[file_process] Không đọc được TMP_DIR để dọn: %s", e)
        return 0

    with entries:
        for entry in entries:
            try:
                if not entry.name.startswith(FILENAME_PREFIX):
                    continue
                if not entry.is_file():
                    continue
                age = now - entry.stat().st_mtime
                if age < FILE_RETENTION_SECONDS:
                    continue
                os.remove(entry.path)
                removed += 1
            except OSError as e:
                logger.warning("[file_process] Lỗi khi dọn '%s': %s", entry.path, e)

    if removed:
        logger.info("[file_process] Cleanup sweep: đã xoá %d file quá 24h.", removed)
    return removed


async def start_cleanup_loop(interval_seconds: int = CLEANUP_SWEEP_INTERVAL_SECONDS) -> None:
    """Vòng lặp nền: định kỳ dọn các file /tmp đã quá FILE_RETENTION_SECONDS.

    Chạy vô hạn cho tới khi bị cancel (main.py cancel task này lúc shutdown).
    Chạy 1 lần ngay lúc start để dọn rác còn sót từ lần restart trước.
    """
    logger.info(
        "[file_process] Cleanup loop bắt đầu — quét mỗi %ds, giữ file %ds (~%.0fh).",
        interval_seconds, FILE_RETENTION_SECONDS, FILE_RETENTION_SECONDS / 3600,
    )
    while True:
        try:
            _sweep_expired_once()
        except Exception as e:
            logger.error("[file_process] Cleanup sweep lỗi bất ngờ: %s", e)
        await asyncio.sleep(interval_seconds)


def _extract_docx(path: str, lang: str = DEFAULT_LANG) -> Optional[str]:
    try:
        import docx  # python-docx
    except ImportError:
        logger.info("[file_process] python-docx chưa được cài — bỏ qua trích xuất .docx")
        return None
    try:
        doc   = docx.Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        text = "\n".join(parts).strip()
    except Exception as e:
        logger.warning("[file_process] Lỗi đọc .docx: %s", e)
        return None
    return _truncate(text, lang) if text else None


def _extract_xlsx(path: str, lang: str = DEFAULT_LANG) -> Optional[str]:
    try:
        import openpyxl
    except ImportError:
        logger.info("[file_process] openpyxl chưa được cài — bỏ qua trích xuất .xlsx")
        return None
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        blocks: list[str] = []
        for ws in wb.worksheets:
            rows_text: list[str] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 200:  # tránh sheet cực lớn nuốt hết budget
                    rows_text.append(t("file.xlsx_rows_truncated", lang))
                    break
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    rows_text.append(" | ".join(cells))
            if rows_text:
                blocks.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows_text))
        text = "\n\n".join(blocks).strip()
    except Exception as e:
        logger.warning("[file_process] Lỗi đọc .xlsx: %s", e)
        return None
    return _truncate(text, lang) if text else None


def _extract_pptx(path: str, lang: str = DEFAULT_LANG) -> Optional[str]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.info("[file_process] python-pptx chưa được cài — bỏ qua trích xuất .pptx")
        return None
    try:
        prs = Presentation(path)
        blocks: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        lines.append(" | ".join(c.text.strip() for c in row.cells))
            if lines:
                blocks.append(f"[Slide {i}]\n" + "\n".join(lines))
        text = "\n\n".join(blocks).strip()
    except Exception as e:
        logger.warning("[file_process] Lỗi đọc .pptx: %s", e)
        return None
    return _truncate(text, lang) if text else None


def _read_text(path: str, lang: str = DEFAULT_LANG) -> Optional[str]:
    try:
        with open(path, "rb") as f:
            raw = f.read()
        text = raw.decode("utf-8", errors="replace").strip()
    except Exception as e:
        logger.warning("[file_process] Lỗi đọc text file: %s", e)
        return None
    return _truncate(text, lang) if text else None


def _truncate(text: str, lang: str = DEFAULT_LANG) -> str:
    if len(text) > MAX_TEXT_CHARS:
        return text[:MAX_TEXT_CHARS] + t("file.truncated", lang)
    return text


async def process_incoming_file(
    bot:      Bot,
    file_id:  str,
    filename: str,
    tg_mime:  Optional[str],
    size:     Optional[int],
    lang:     str = DEFAULT_LANG,
) -> FileResult:
    """Tải file thật + trả về nội dung dùng được cho Gemini (nếu có).

    Lưu ý: file tải về TMP_DIR KHÔNG bị xoá ở cuối hàm này nữa — được giữ
    lại và dọn định kỳ bởi start_cleanup_loop() (xoá sau FILE_RETENTION_SECONDS).
    """
    path = await _download(bot, file_id, filename, size)
    if not path:
        return FileResult()  # quá lớn hoặc tải lỗi → fallback metadata-only

    mime = _guess_mime(filename, tg_mime)
    ext  = os.path.splitext(filename)[1].lower()

    if mime in GEMINI_INLINE_MIMES:
        with open(path, "rb") as f:
            data = f.read()
        b64 = base64.b64encode(data).decode("ascii")
        return FileResult(media_part={"inlineData": {"mimeType": mime, "data": b64}})

    if ext == ".docx":
        text = _extract_docx(path, lang)
        return FileResult(extra_text=text) if text else FileResult()

    if ext == ".xlsx":
        text = _extract_xlsx(path, lang)
        return FileResult(extra_text=text) if text else FileResult()

    if ext == ".pptx":
        text = _extract_pptx(path, lang)
        return FileResult(extra_text=text) if text else FileResult()

    if ext in TEXT_EXTS or mime.startswith("text/"):
        text = _read_text(path, lang)
        return FileResult(extra_text=text) if text else FileResult()

    # Định dạng chưa hỗ trợ trích xuất (vd .zip, .xls cũ, .rar...)
    return FileResult()
